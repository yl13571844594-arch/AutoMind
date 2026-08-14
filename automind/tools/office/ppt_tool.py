"""PowerPoint 工具 —— 生成 / 读取 / 追加幻灯片（.pptx）。

社区版动作：create / read / add_slide
专业版动作（office_pro）：template（母版/主题套用）
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from automind.core.types import PermissionTier, ToolResult
from automind.tools._toolkit import bad, delegate_pro, err, need, ok
from automind.tools.base import AbstractTool

#: 进阶动作（由专业版 office_pro 实现，社区版仅转交）
PRO_ACTIONS = {"template"}


class PptTool(AbstractTool):
    """生成与读取 PowerPoint 演示文稿（.pptx）。"""

    name = "ppt_tool"
    description = (
        "Create, read and extend PowerPoint decks (.pptx). Actions: "
        "create (build a new deck from a list of slides, each with a title and bullets), "
        "read (extract text of every slide), "
        "add_slide (append a titled slide with bullets to an existing deck). "
        "The template action (apply a master/theme) requires the Pro edition."
    )
    parameters = {
        "type": "object",
        "properties": {
            "action": {"type": "string", "enum": ["create", "read", "add_slide", "template"]},
            "path": {"type": "string", "description": "Deck path (.pptx)."},
            "title": {"type": "string", "description": "Slide title (for add_slide)."},
            "bullets": {
                "type": "array", "items": {"type": "string"},
                "description": "Bullet points for the slide (create/add_slide).",
            },
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                    },
                },
                "description": "Slides for create: list of {title, bullets}.",
            },
        },
        "required": ["action", "path"],
    }
    permission_tier = PermissionTier.SENSITIVE
    risk_score = 30

    async def execute(self, **kwargs: Any) -> ToolResult:
        action = str(kwargs.get("action", "")).lower()
        path = Path(str(kwargs.get("path", ""))).expanduser()
        try:
            if action in PRO_ACTIONS:
                need("pptx")
                return ok(self.name, **delegate_pro("office_pro", self.name, action, kwargs))
            need("pptx")
            from pptx import Presentation
            from pptx.util import Pt

            if action == "create":
                slides = kwargs.get("slides") or []
                if not isinstance(slides, list):
                    slides = []
                if not slides and (kwargs.get("title") or kwargs.get("bullets")):
                    # 顶层同时存在 title/bullets（add_slide 的参数形状），模型很容易
                    # 拿它们来调 create。与其回一句"参数不对"让它重试一轮，
                    # 不如按"单页 PPT"理解 —— 意图是明确的。
                    slides = [{"title": kwargs.get("title", ""),
                               "bullets": kwargs.get("bullets") or []}]
                if not slides:
                    return bad(self.name,
                               "create 需要提供 slides 列表（每项含 title 与 bullets）；"
                               "只做一页时也可直接传 title 与 bullets")
                prs = Presentation()
                for slide in slides:
                    title = str(slide.get("title", "")) if isinstance(slide, dict) else str(slide)
                    bullets = slide.get("bullets", []) if isinstance(slide, dict) else []
                    _append_slide(prs, title, [str(b) for b in bullets], Pt)
                path.parent.mkdir(parents=True, exist_ok=True)
                prs.save(str(path))
                return ok(self.name, action=action, path=str(path),
                          slides=len(slides), message=f"已生成 {len(slides)} 页 PPT：{path}")

            if action == "add_slide":
                if not path.exists():
                    return bad(self.name, f"文件不存在：{path}（可先用 create 新建）")
                prs = Presentation(str(path))
                title = str(kwargs.get("title", ""))
                bullets = [str(b) for b in (kwargs.get("bullets") or [])]
                _append_slide(prs, title, bullets, Pt)
                prs.save(str(path))
                return ok(self.name, action=action, path=str(path),
                          total_slides=len(prs.slides), message="已追加一页")

            if action == "read":
                if not path.exists():
                    return bad(self.name, f"文件不存在：{path}")
                prs = Presentation(str(path))
                slides_out = []
                for i, slide in enumerate(prs.slides, 1):
                    texts = []
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False):
                            for para in shape.text_frame.paragraphs:
                                line = "".join(r.text for r in para.runs).strip()
                                if line:
                                    texts.append(line)
                    slides_out.append({"slide": i, "text": "\n".join(texts)})
                return ok(self.name, action=action, path=str(path), slides=slides_out)

            return bad(self.name, f"未知 action：{action}")
        except Exception as e:
            return err(self.name, e)


def _append_slide(prs: Any, title: str, bullets: list[str], pt: Any) -> None:
    """向演示文稿追加一页「标题 + 要点」。

    Args:
        pt: python-pptx 的 ``Pt`` 长度构造器（由调用方传入，避免重复导入）。
    """
    layout = prs.slide_layouts[1]  # 标题和内容版式
    slide = prs.slides.add_slide(layout)
    if title:
        slide.shapes.title.text = title
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
    if body is not None and body.has_text_frame:
        tf = body.text_frame
        tf.clear()
        for i, line in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.level = 0
            for run in para.runs:
                run.font.size = pt(18)
    return
